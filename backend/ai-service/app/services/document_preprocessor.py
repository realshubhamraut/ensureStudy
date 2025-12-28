"""
Document Preprocessor - Multi-Format Document Handler

Supports:
- PDF: Extract text and images, detect scanned vs text-based
- PPTX: Convert slides to images
- DOCX: Extract text and embedded images
- Images: Deskew, contrast enhancement, noise reduction

All processing is done locally with OpenCV and standard Python libraries.
"""
import logging
import os
import tempfile
from typing import Dict, Any, List, Optional
from pathlib import Path
import asyncio

import cv2
import numpy as np

logger = logging.getLogger(__name__)


class DocumentPreprocessor:
    """
    Multi-format document preprocessor
    
    Converts various document formats into processed images
    ready for OCR.
    """
    
    def __init__(self, output_dpi: int = 200):
        self.output_dpi = output_dpi
        # A4 size at output DPI
        self.page_width = int(8.27 * output_dpi)
        self.page_height = int(11.69 * output_dpi)
    
    async def process_pdf(self, pdf_path: str) -> Dict[str, Any]:
        """
        Process PDF file
        
        Returns:
            {image_paths: [...], text_content: str, is_scanned: bool}
        """
        try:
            import pypdf
        except ImportError:
            return await self._process_pdf_with_images(pdf_path)
        
        output_dir = os.path.join(os.path.dirname(pdf_path), "preprocessed")
        os.makedirs(output_dir, exist_ok=True)
        
        image_paths = []
        text_content = ""
        is_scanned = False
        
        try:
            reader = pypdf.PdfReader(pdf_path)
            
            for i, page in enumerate(reader.pages):
                # Extract text
                page_text = page.extract_text() or ""
                text_content += page_text + "\n"
                
                # Check if this is a scanned page (little text)
                if len(page_text.strip()) < 100:
                    is_scanned = True
            
            # If scanned or has images, convert to images
            if is_scanned or len(text_content.strip()) < 500:
                # Use pdf2image or fallback
                image_paths = await self._pdf_to_images(pdf_path, output_dir)
            else:
                # Text-based PDF - still create images for consistent processing
                image_paths = await self._pdf_to_images(pdf_path, output_dir)
            
            return {
                "image_paths": image_paths,
                "text_content": text_content,
                "is_scanned": is_scanned,
                "page_count": len(reader.pages)
            }
            
        except Exception as e:
            logger.error(f"PDF processing error: {e}")
            raise
    
    async def _pdf_to_images(self, pdf_path: str, output_dir: str) -> List[str]:
        """Convert PDF pages to images"""
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(
                pdf_path,
                dpi=self.output_dpi,
                fmt='png'
            )
            
            image_paths = []
            for i, img in enumerate(images):
                output_path = os.path.join(output_dir, f"page_{i+1:03d}.png")
                img.save(output_path, 'PNG')
                image_paths.append(output_path)
            
            return image_paths
            
        except ImportError:
            logger.warning("pdf2image not available, using fallback")
            return await self._process_pdf_with_images(pdf_path)
    
    async def _process_pdf_with_images(self, pdf_path: str) -> Dict[str, Any]:
        """Fallback PDF processing - just return metadata"""
        return {
            "image_paths": [],
            "text_content": "",
            "is_scanned": True,
            "error": "pdf2image not installed"
        }
    
    async def process_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """
        Process PowerPoint file
        
        Extracts slides as images and speaker notes.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            logger.error("python-pptx not installed")
            return {"image_paths": [], "error": "python-pptx not installed"}
        
        output_dir = os.path.join(os.path.dirname(pptx_path), "preprocessed")
        os.makedirs(output_dir, exist_ok=True)
        
        try:
            prs = Presentation(pptx_path)
            
            image_paths = []
            speaker_notes = []
            
            for i, slide in enumerate(prs.slides):
                # Extract speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    speaker_notes.append(notes)
                else:
                    speaker_notes.append("")
                
                # For now, we can't easily render slides to images without
                # additional dependencies. We'll extract text and shapes instead.
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                # Create a placeholder image with text
                output_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
                self._create_text_image(
                    "\n".join(slide_text),
                    output_path,
                    title=f"Slide {i+1}"
                )
                image_paths.append(output_path)
            
            return {
                "image_paths": image_paths,
                "speaker_notes": speaker_notes,
                "slide_count": len(prs.slides)
            }
            
        except Exception as e:
            logger.error(f"PPTX processing error: {e}")
            raise
    
    async def process_docx(self, docx_path: str) -> Dict[str, Any]:
        """
        Process Word document
        
        Extracts text and embedded images.
        """
        try:
            import mammoth
        except ImportError:
            logger.error("mammoth not installed")
            return {"image_paths": [], "error": "mammoth not installed"}
        
        output_dir = os.path.join(os.path.dirname(docx_path), "preprocessed")
        os.makedirs(output_dir, exist_ok=True)
        
        try:
            with open(docx_path, "rb") as f:
                result = mammoth.convert_to_html(f)
            
            html_content = result.value
            text_content = mammoth.extract_raw_text(open(docx_path, "rb")).value
            
            # Create images from text (paginated)
            image_paths = self._text_to_images(text_content, output_dir)
            
            return {
                "image_paths": image_paths,
                "text_content": text_content,
                "html_content": html_content
            }
            
        except Exception as e:
            logger.error(f"DOCX processing error: {e}")
            raise
    
    async def process_images(self, image_paths: List[str]) -> Dict[str, Any]:
        """
        Process image files
        
        Applies:
        - Deskew correction
        - Contrast enhancement
        - Noise reduction
        """
        output_dir = os.path.join(os.path.dirname(image_paths[0]), "preprocessed")
        os.makedirs(output_dir, exist_ok=True)
        
        processed_paths = []
        
        for i, img_path in enumerate(image_paths):
            try:
                img = cv2.imread(img_path)
                if img is None:
                    logger.warning(f"Could not read image: {img_path}")
                    continue
                
                # Apply preprocessing
                processed = self._preprocess_image(img)
                
                # Save
                output_path = os.path.join(output_dir, f"page_{i+1:03d}.png")
                cv2.imwrite(output_path, processed)
                processed_paths.append(output_path)
                
            except Exception as e:
                logger.warning(f"Image processing error for {img_path}: {e}")
                # Copy original as fallback
                import shutil
                output_path = os.path.join(output_dir, f"page_{i+1:03d}.png")
                shutil.copy(img_path, output_path)
                processed_paths.append(output_path)
        
        return {
            "image_paths": processed_paths,
            "page_count": len(processed_paths)
        }
    
    def _preprocess_image(self, img: np.ndarray) -> np.ndarray:
        """
        Apply image preprocessing for OCR
        
        Steps:
        1. Resize if too large
        2. Deskew
        3. Contrast enhancement (CLAHE)
        4. Noise reduction
        """
        # 1. Resize if needed
        h, w = img.shape[:2]
        max_dim = max(self.page_width, self.page_height)
        
        if max(h, w) > max_dim:
            scale = max_dim / max(h, w)
            img = cv2.resize(img, None, fx=scale, fy=scale, interpolation=cv2.INTER_AREA)
        
        # 2. Convert to grayscale for processing
        if len(img.shape) == 3:
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        else:
            gray = img
        
        # 3. Deskew (simple approach)
        gray = self._deskew(gray)
        
        # 4. Contrast enhancement with CLAHE
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(gray)
        
        # 5. Light noise reduction
        denoised = cv2.fastNlMeansDenoising(enhanced, None, h=10)
        
        # Convert back to BGR for consistency
        return cv2.cvtColor(denoised, cv2.COLOR_GRAY2BGR)
    
    def _deskew(self, img: np.ndarray, max_angle: float = 10.0) -> np.ndarray:
        """Deskew image using Hough transform"""
        try:
            # Edge detection
            edges = cv2.Canny(img, 50, 150, apertureSize=3)
            
            # Hough lines
            lines = cv2.HoughLines(edges, 1, np.pi/180, 200)
            
            if lines is None:
                return img
            
            # Calculate dominant angle
            angles = []
            for line in lines[:20]:  # Use first 20 lines
                rho, theta = line[0]
                angle = (theta * 180 / np.pi) - 90
                if abs(angle) < max_angle:
                    angles.append(angle)
            
            if not angles:
                return img
            
            median_angle = np.median(angles)
            
            # Rotate image
            h, w = img.shape[:2]
            center = (w // 2, h // 2)
            M = cv2.getRotationMatrix2D(center, median_angle, 1.0)
            rotated = cv2.warpAffine(img, M, (w, h), 
                                      borderMode=cv2.BORDER_REPLICATE)
            
            return rotated
            
        except Exception as e:
            logger.warning(f"Deskew error: {e}")
            return img
    
    def _create_text_image(
        self,
        text: str,
        output_path: str,
        title: str = ""
    ) -> None:
        """Create an image from text (for PPTX slides without rendering)"""
        # Create white image
        img = np.ones((self.page_height, self.page_width, 3), dtype=np.uint8) * 255
        
        # Add title
        if title:
            cv2.putText(img, title, (50, 80), 
                       cv2.FONT_HERSHEY_SIMPLEX, 2.0, (0, 0, 0), 3)
        
        # Add text (wrapped)
        y = 150
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.8
        
        for line in text.split('\n')[:30]:  # Limit lines
            if y > self.page_height - 50:
                break
            cv2.putText(img, line[:100], (50, y), font, font_scale, (0, 0, 0), 1)
            y += 40
        
        cv2.imwrite(output_path, img)
    
    def _text_to_images(self, text: str, output_dir: str) -> List[str]:
        """Convert long text to multiple page images"""
        lines = text.split('\n')
        lines_per_page = 40
        
        image_paths = []
        
        for page_num, i in enumerate(range(0, len(lines), lines_per_page)):
            page_lines = lines[i:i + lines_per_page]
            page_text = '\n'.join(page_lines)
            
            output_path = os.path.join(output_dir, f"page_{page_num+1:03d}.png")
            self._create_text_image(page_text, output_path)
            image_paths.append(output_path)
        
        return image_paths


# Convenience function
async def preprocess_document(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    Preprocess a document file
    
    Args:
        file_path: Path to the document
        file_type: File extension (pdf, pptx, docx, jpg, etc.)
    
    Returns:
        {image_paths: [...], ...}
    """
    preprocessor = DocumentPreprocessor()
    
    file_type = file_type.lower().strip('.')
    
    if file_type == "pdf":
        return await preprocessor.process_pdf(file_path)
    elif file_type == "pptx":
        return await preprocessor.process_pptx(file_path)
    elif file_type == "docx":
        return await preprocessor.process_docx(file_path)
    elif file_type in ["jpg", "jpeg", "png", "webp", "heic"]:
        return await preprocessor.process_images([file_path])
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
