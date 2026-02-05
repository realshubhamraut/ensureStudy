"""
PPTX Extractor Service - Extract text content from PowerPoint presentations

Uses python-pptx to parse PPTX files and extract:
- Slide text content
- Slide count
- Image detection

Also provides PPTX to PDF conversion using LibreOffice headless.
"""
import os
import logging
import subprocess
import shutil
from typing import Tuple, Optional
from dataclasses import dataclass

logger = logging.getLogger(__name__)


def convert_pptx_to_pdf(pptx_path: str, output_dir: Optional[str] = None) -> Optional[str]:
    """
    Convert PPTX to PDF using LibreOffice headless mode.
    
    Args:
        pptx_path: Path to the PPTX file
        output_dir: Directory for output PDF (defaults to same as PPTX)
        
    Returns:
        Path to the generated PDF file, or None if conversion failed
    """
    if not os.path.exists(pptx_path):
        logger.error(f"[PPTX→PDF] File not found: {pptx_path}")
        return None
    
    if output_dir is None:
        output_dir = os.path.dirname(pptx_path)
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Find LibreOffice executable
    soffice_paths = [
        'soffice',  # Linux/macOS if in PATH
        '/usr/bin/soffice',  # Linux
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS
        'C:\\Program Files\\LibreOffice\\program\\soffice.exe',  # Windows
    ]
    
    soffice_cmd = None
    for path in soffice_paths:
        if shutil.which(path) or os.path.exists(path):
            soffice_cmd = path
            break
    
    if not soffice_cmd:
        logger.error("[PPTX→PDF] LibreOffice not found. Install: brew install --cask libreoffice")
        return None
    
    try:
        print(f"[PPTX→PDF] Converting: {os.path.basename(pptx_path)}")
        
        result = subprocess.run(
            [soffice_cmd, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
            capture_output=True,
            text=True,
            timeout=120  # 2 minute timeout for large files
        )
        
        if result.returncode != 0:
            logger.error(f"[PPTX→PDF] Conversion failed: {result.stderr}")
            return None
        
        # Get the output PDF path
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
        
        if os.path.exists(pdf_path):
            print(f"[PPTX→PDF] ✅ Converted: {pdf_path}")
            return pdf_path
        else:
            logger.error(f"[PPTX→PDF] PDF not created at expected path: {pdf_path}")
            return None
            
    except subprocess.TimeoutExpired:
        logger.error("[PPTX→PDF] Conversion timed out (>120s)")
        return None
    except Exception as e:
        logger.error(f"[PPTX→PDF] Error: {e}")
        return None


@dataclass
class PPTXExtractionResult:
    """Result of PPTX text extraction."""
    text: str
    slide_count: int
    has_images: bool
    error: Optional[str] = None


class PPTXExtractor:
    """Extract text content from PowerPoint presentations."""
    
    def extract_text_from_pptx(self, file_path: str) -> Tuple[str, bool, int]:
        """
        Extract text from all slides in a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Tuple of (text_content, has_images, slide_count)
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            if not os.path.exists(file_path):
                logger.error(f"[PPTX] File not found: {file_path}")
                return ("", False, 0)
            
            print(f"[PPTX] Extracting text from: {os.path.basename(file_path)}")
            
            prs = Presentation(file_path)
            
            all_text = []
            has_images = False
            slide_count = len(prs.slides)
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n--- Slide {slide_idx} ---\n"]
                
                for shape in slide.shapes:
                    # Check for images
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_images = True
                    
                    # Extract text from shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the header
                    all_text.append("\n".join(slide_text))
            
            full_text = "\n\n".join(all_text)
            
            print(f"[PPTX] ✅ Extracted {len(full_text)} chars from {slide_count} slides")
            
            return (full_text, has_images, slide_count)
            
        except ImportError:
            logger.error("[PPTX] python-pptx not installed. Run: pip install python-pptx")
            return ("", False, 0)
        except Exception as e:
            logger.error(f"[PPTX] Extraction error: {e}")
            return ("", False, 0)
    
    def extract(self, file_path: str) -> PPTXExtractionResult:
        """
        Extract text from PPTX and return structured result.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            PPTXExtractionResult with text, slide_count, has_images
        """
        try:
            text, has_images, slide_count = self.extract_text_from_pptx(file_path)
            return PPTXExtractionResult(
                text=text,
                slide_count=slide_count,
                has_images=has_images
            )
        except Exception as e:
            return PPTXExtractionResult(
                text="",
                slide_count=0,
                has_images=False,
                error=str(e)
            )


# Singleton instance
pptx_extractor = PPTXExtractor()


# ============================================================================
# Test
# ============================================================================

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        result = pptx_extractor.extract(file_path)
        print(f"\n{'='*60}")
        print(f"Slides: {result.slide_count}")
        print(f"Has Images: {result.has_images}")
        print(f"Text Length: {len(result.text)} chars")
        print(f"{'='*60}")
        print(result.text[:2000] if result.text else "No text extracted")
    else:
        print("Usage: python pptx_extractor.py <path_to_pptx>")
